from __future__ import annotations

import csv
import hashlib
import json
from collections.abc import Iterable
from io import StringIO
from pathlib import Path

PLAIN_TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".py",
    ".pyi",
    ".js",
    ".jsx",
    ".mjs",
    ".cjs",
    ".ts",
    ".tsx",
    ".mts",
    ".cts",
    ".java",
    ".kt",
    ".kts",
    ".go",
    ".rs",
    ".c",
    ".h",
    ".cc",
    ".cpp",
    ".cxx",
    ".hh",
    ".hpp",
    ".hxx",
    ".cs",
    ".fs",
    ".fsx",
    ".rb",
    ".php",
    ".swift",
    ".scala",
    ".sh",
    ".bash",
    ".zsh",
    ".fish",
    ".ps1",
    ".sql",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".conf",
    ".xml",
    ".vue",
    ".svelte",
    ".proto",
    ".graphql",
    ".gql",
    ".dml",
    ".simics",
    ".mk",
    ".inc",
    ".include",
    ".cmake",
}
PLAIN_TEXT_FILENAMES = {
    "dockerfile",
    "makefile",
    "gnumakefile",
    "rakefile",
    "gemfile",
    "procfile",
    "jenkinsfile",
}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".tif", ".tiff", ".webp"}
SUPPORTED_EXTENSIONS = (
    PLAIN_TEXT_EXTENSIONS
    | {".json", ".csv", ".html", ".htm", ".pdf", ".docx", ".pptx"}
    | IMAGE_EXTENSIONS
)

# Embedded images repeat on every page of a document as logos and stamps. Only the
# first occurrence of a given image is worth reading.
MIN_EMBEDDED_IMAGE_BYTES = 3000
MIN_OCR_CONFIDENCE = 0.5

_ocr_reader: object | None = None
_ocr_import_failed = False


def _ocr_engine() -> object | None:
    """Return a shared RapidOCR engine, or None when the package is unavailable."""
    global _ocr_reader, _ocr_import_failed
    if _ocr_import_failed:
        return None
    if _ocr_reader is None:
        try:
            from rapidocr_onnxruntime import RapidOCR
        except ImportError:
            _ocr_import_failed = True
            return None
        _ocr_reader = RapidOCR()
    return _ocr_reader


def ocr_text(image: bytes | str | Path) -> str:
    """Read the text in an image. Returns an empty string when OCR is unavailable."""
    engine = _ocr_engine()
    if engine is None:
        return ""
    try:
        result, _ = engine(image if isinstance(image, bytes) else str(image))
    except Exception:
        return ""
    if not result:
        return ""
    return " ".join(
        text.strip()
        for _, text, score in result
        if text and text.strip() and score >= MIN_OCR_CONFIDENCE
    )


def _embedded_image_text(blobs: Iterable[bytes], seen_hashes: set[str]) -> str:
    """OCR the images of one page or document, skipping repeats and tiny icons."""
    if _ocr_engine() is None:
        return ""
    texts: list[str] = []
    for data in blobs:
        if len(data) < MIN_EMBEDDED_IMAGE_BYTES:
            continue
        digest = hashlib.md5(data).hexdigest()
        if digest in seen_hashes:
            continue
        seen_hashes.add(digest)
        text = ocr_text(data)
        if text:
            texts.append(text)
    return "\n".join(texts)


def parse_file(path: Path) -> str:
    extension = path.suffix.lower()
    if extension in PLAIN_TEXT_EXTENSIONS or path.name.lower() in PLAIN_TEXT_FILENAMES:
        return path.read_text(encoding="utf-8", errors="replace")
    if extension == ".json":
        data = json.loads(path.read_text(encoding="utf-8"))
        return json.dumps(data, ensure_ascii=False, indent=2)
    if extension == ".csv":
        content = path.read_text(encoding="utf-8", errors="replace")
        reader = csv.reader(StringIO(content))
        return "\n".join(" | ".join(cell.strip() for cell in row) for row in reader)
    if extension in {".html", ".htm"}:
        return _parse_html(path)
    if extension == ".pdf":
        return _parse_pdf(path)
    if extension == ".docx":
        return _parse_docx(path)
    if extension == ".pptx":
        return _parse_pptx(path)
    if extension in IMAGE_EXTENSIONS:
        return _parse_image(path)
    raise ValueError(f"Unsupported file extension: {extension}")


def _parse_image(path: Path) -> str:
    if _ocr_engine() is None:
        raise RuntimeError(
            "rapidocr-onnxruntime is required to ingest image files. Run: pnpm rag:install"
        )
    return ocr_text(path)


def _parse_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise RuntimeError("beautifulsoup4 is required to ingest HTML files. Run: pnpm rag:install") from exc

    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"), "html.parser")
    for element in soup(["script", "style", "noscript"]):
        element.decompose()
    return soup.get_text(separator="\n", strip=True)


def _parse_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf is required to ingest PDF files. Run: pnpm rag:install") from exc

    reader = PdfReader(str(path))
    pages: list[str] = []
    seen_hashes: set[str] = set()
    for index, page in enumerate(reader.pages):
        parts: list[str] = []
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())
        image_text = _embedded_image_text(_pdf_page_blobs(page), seen_hashes)
        if image_text:
            parts.append(f"(image) {image_text}")
        if parts:
            pages.append(f"[page {index + 1}]\n" + "\n".join(parts))
    return "\n\n".join(pages)


def _pdf_page_blobs(page: object) -> Iterable[bytes]:
    try:
        images = list(page.images)
    except Exception:
        return []
    blobs: list[bytes] = []
    for image in images:
        try:
            blobs.append(image.data)
        except Exception:
            continue
    return blobs


def _parse_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required to ingest DOCX files. Run: pnpm rag:install") from exc

    document = Document(str(path))
    parts: list[str] = []
    parts.extend(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    image_text = _embedded_image_text(_docx_blobs(document), set())
    if image_text:
        parts.append(f"(image) {image_text}")
    return "\n".join(parts)


def _docx_blobs(document: object) -> Iterable[bytes]:
    try:
        image_parts = list(document.part.package.image_parts)
    except Exception:
        return []
    blobs: list[bytes] = []
    for image_part in image_parts:
        try:
            blobs.append(image_part.blob)
        except Exception:
            continue
    return blobs


def _parse_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to ingest PPTX files. Run: pnpm rag:install") from exc

    presentation = Presentation(str(path))
    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            slides.append(f"[slide {slide_index + 1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)
